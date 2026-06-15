"""Build the executive .pptx deck from rev 3 of the markdown source.

Run from repo root:
    python proposals/_build_deck.py

Produces:
    proposals/exec-deck-anthropic-self-service-analytics-2026-06-04.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


# ──────────────────────────── theme ────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
INK = RGBColor(0x1A, 0x1A, 0x2E)
SUBTLE = RGBColor(0x55, 0x66, 0x77)
ACCENT = RGBColor(0x0E, 0x6B, 0xA8)
GREEN_FILL = RGBColor(0xC8, 0xE6, 0xC9)
YELLOW_FILL = RGBColor(0xFF, 0xF3, 0xC4)
RED_FILL = RGBColor(0xFF, 0xCD, 0xD2)
GREEN_DOT = RGBColor(0x2E, 0x7D, 0x32)
YELLOW_DOT = RGBColor(0xF9, 0xA8, 0x25)
RED_DOT = RGBColor(0xC6, 0x28, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_FILL = RGBColor(0xEC, 0xEF, 0xF4)
BAND_FILL = RGBColor(0xF7, 0xF9, 0xFC)
RULE = RGBColor(0xDB, 0xDF, 0xE6)


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def add_text(slide, left, top, width, height, text, *, size=14, bold=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(0.04)
    box.text_frame.margin_right = Inches(0.04)
    box.text_frame.margin_top = Inches(0.02)
    box.text_frame.margin_bottom = Inches(0.02)
    box.text_frame.vertical_anchor = anchor
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Segoe UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_rich(slide, left, top, width, height, runs, *, anchor=MSO_ANCHOR.TOP,
             align=PP_ALIGN.LEFT, line_spacing=1.15):
    """runs = list of dicts {'text': str, 'size': int, 'bold': bool, 'color': rgb, 'italic': bool}"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    for i, item in enumerate(runs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if "space_before" in item:
            p.space_before = Pt(item["space_before"])
        if "space_after" in item:
            p.space_after = Pt(item["space_after"])
        for r in item["runs"]:
            run = p.add_run()
            run.text = r["text"]
            run.font.name = r.get("font", "Segoe UI")
            run.font.size = Pt(r.get("size", 14))
            run.font.bold = r.get("bold", False)
            run.font.italic = r.get("italic", False)
            run.font.color.rgb = r.get("color", INK)
    return box


def add_rect(slide, left, top, width, height, fill_color, line_color=None,
             line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    shape.shadow.inherit = False
    return shape


def slide_header(slide, title, *, eyebrow=None, page_num=None, total=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.07), NAVY)
    if eyebrow:
        add_text(slide, Inches(0.5), Inches(0.20), Inches(8), Inches(0.3),
                 eyebrow.upper(), size=10, bold=True, color=ACCENT)
    add_text(slide, Inches(0.5), Inches(0.45), Inches(12), Inches(0.7), title,
             size=26, bold=True, color=NAVY)
    if page_num is not None and total is not None:
        add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
                 f"{page_num} / {total}", size=9, color=SUBTLE,
                 align=PP_ALIGN.RIGHT)
    # Footer hairline
    add_rect(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.01), RULE)


def add_bullet_list(slide, left, top, width, height, items, *, size=14,
                    spacing=8, indent_size=12):
    runs = []
    for i, item in enumerate(items):
        if isinstance(item, str):
            runs.append({
                "runs": [
                    {"text": "•  ", "size": size, "bold": True, "color": ACCENT},
                    {"text": item, "size": size, "color": INK},
                ],
                "space_before": 0 if i == 0 else spacing,
            })
        else:
            # tuple: (lead, body) — first run is bold lead
            lead, body = item
            runs.append({
                "runs": [
                    {"text": "•  ", "size": size, "bold": True, "color": ACCENT},
                    {"text": lead, "size": size, "bold": True, "color": INK},
                    {"text": body, "size": size, "color": INK},
                ],
                "space_before": 0 if i == 0 else spacing,
            })
    return add_rich(slide, left, top, width, height, runs)


# ──────────────────────────── slide builders ────────────────────────────

def build_slide_1_cover(prs, total):
    s = add_blank_slide(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, Inches(0), Inches(5.0), SLIDE_W, Inches(2.5), RGBColor(0x06, 0x1A, 0x2E))
    add_text(s, Inches(1.0), Inches(1.4), Inches(11), Inches(0.4),
             "EXECUTIVE BRIEFING · 2026-06-04", size=12, bold=True,
             color=RGBColor(0x6F, 0xB0, 0xD9))
    add_text(s, Inches(1.0), Inches(1.95), Inches(11.3), Inches(1.5),
             "From static corpus to living business brain",
             size=40, bold=True, color=WHITE)
    add_text(s, Inches(1.0), Inches(3.30), Inches(11.3), Inches(1.4),
             "Self-service analytics at eToro — where we stand vs Anthropic, "
             "and how we build the machine",
             size=20, color=RGBColor(0xB7, 0xCF, 0xE6))
    add_text(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(0.5),
             "Source: How Anthropic enables self-service data analytics with Claude (claude.com/blog, 2026-06-03)",
             size=12, color=RGBColor(0xB7, 0xCF, 0xE6), italic=True)
    add_text(s, Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.5),
             "Audience: Data leadership · Eng leadership · Analyst-team leads   |   Read time: ~10 min",
             size=12, color=RGBColor(0xB7, 0xCF, 0xE6))


def build_slide_2_why(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "Why this matters", eyebrow="Context",
                 page_num=page, total=total)
    add_bullet_list(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5), [
        ("Anthropic's published number: ",
         "95% of business analytics questions answered automatically, ~95% accurate."),
        ("Same primitives we have: ",
         "Claude, MCP, skills, dbt-style modeling, an emerging metric layer."),
        ("We've built ~80% of the stack already — ",
         "including two assets most teams would have to build from zero."),
        ("The strategic shift for H2: ",
         "stop treating the skill corpus as a static artifact. Start running it as a living machine — one that watches every knowledge source in the company, judges itself with telemetry, and updates itself."),
        ("The ask: ",
         "align on the machine we're building + commit to a measurable accuracy floor."),
    ], size=17, spacing=16)


def build_slide_3_framework(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "Anthropic's framework", eyebrow="The model we'll keep referring to",
                 page_num=page, total=total)
    # Pull quote
    add_rect(s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.7), HEADER_FILL)
    add_text(s, Inches(0.95), Inches(1.4), Inches(11.6), Inches(0.55),
             "Self-service analytics accuracy is a context + verification problem — not a code-generation problem.",
             size=15, italic=True, color=INK)

    # Three failure modes
    add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(0.35),
             "Three failure modes that cause most wrong answers",
             size=14, bold=True, color=NAVY)
    fm_data = [
        ("Concept ↔ entity ambiguity",
         "agent can't map \"revenue\" / \"active user\" / \"FTD\" to the one correct table + filter."),
        ("Data staleness",
         "schemas + definitions drift. Anthropic measured 95% → 65% in 1 month without active maintenance."),
        ("Retrieval failure",
         "right info is in the docs, agent doesn't find it. Grep over thousands of notebooks moved accuracy <1pt. Structure > volume."),
    ]
    col_w = Inches(4.0)
    for i, (h, body) in enumerate(fm_data):
        left = Inches(0.7 + i * 4.1)
        add_rect(s, left, Inches(2.65), col_w, Inches(1.45), BAND_FILL,
                 line_color=RULE)
        add_text(s, left + Inches(0.15), Inches(2.75), col_w - Inches(0.3),
                 Inches(0.4), h, size=12, bold=True, color=ACCENT)
        add_text(s, left + Inches(0.15), Inches(3.10), col_w - Inches(0.3),
                 Inches(1.05), body, size=11, color=INK)

    # Four-layer stack
    add_text(s, Inches(0.7), Inches(4.35), Inches(12), Inches(0.35),
             "Four-layer stack that attacks each",
             size=14, bold=True, color=NAVY)
    layers = [
        ("Data foundations", "Canonical datasets, enforced governance, colocation."),
        ("Sources of truth", "Semantic layer FIRST → lineage → query corpus → business context."),
        ("Skills", "Pairwise: knowledge skill (router) + unbook skill (process + analysis patterns)."),
        ("Validation", "Offline evals (pinned) + ablations + online provenance + correction harvesting."),
    ]
    row_top = Inches(4.75)
    row_h = Inches(0.45)
    for i, (h, body) in enumerate(layers):
        top = row_top + i * row_h
        add_rect(s, Inches(0.7), top, Inches(12), row_h, BAND_FILL if i % 2 == 0 else WHITE,
                 line_color=RULE)
        add_text(s, Inches(0.9), top + Inches(0.05), Inches(3.3), row_h, h,
                 size=12, bold=True, color=NAVY)
        add_text(s, Inches(4.2), top + Inches(0.05), Inches(8.4), row_h, body,
                 size=12, color=INK)


def add_scorecard_table(slide, rows, left, top, total_w, total_h):
    """rows = list of (status_dot, layer, status_text, priority).
    status_dot ∈ {'g','y','r'}.
    """
    col_widths = [Inches(0.5), Inches(3.3), Inches(7.3), Inches(1.1)]
    header_h = Inches(0.36)
    body_top = top + header_h
    available_h = total_h - header_h
    row_h = available_h / len(rows) if rows else Inches(0.3)

    # Header band
    add_rect(slide, left, top, total_w, header_h, NAVY)
    x = left
    headers = ["", "Anthropic layer", "eToro today", "Priority"]
    for w, h in zip(col_widths, headers):
        add_text(slide, x + Inches(0.1), top + Inches(0.05), w - Inches(0.2),
                 header_h - Inches(0.1), h, size=11, bold=True, color=WHITE)
        x += w

    dot_color_map = {"g": GREEN_DOT, "y": YELLOW_DOT, "r": RED_DOT}
    fill_map = {"g": GREEN_FILL, "y": YELLOW_FILL, "r": RED_FILL}

    for i, (dot, layer, status, prio) in enumerate(rows):
        row_top = body_top + int(row_h * i)
        # Row background banding
        add_rect(slide, left, row_top, total_w, row_h,
                 BAND_FILL if i % 2 == 0 else WHITE, line_color=RULE)
        # Status pill cell
        pill_left = left + Inches(0.10)
        pill_top = row_top + Inches(0.05)
        pill_w = Inches(0.3)
        pill_h = row_h - Inches(0.10)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, pill_left, pill_top,
                                      pill_w, pill_h)
        circ.fill.solid()
        circ.fill.fore_color.rgb = dot_color_map[dot]
        circ.line.fill.background()
        circ.shadow.inherit = False
        # Layer name
        add_text(slide, left + col_widths[0] + Inches(0.05),
                 row_top + Inches(0.02),
                 col_widths[1] - Inches(0.1), row_h - Inches(0.04),
                 layer, size=10, bold=True, color=INK,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Status text
        add_text(slide, left + col_widths[0] + col_widths[1] + Inches(0.05),
                 row_top + Inches(0.02),
                 col_widths[2] - Inches(0.1), row_h - Inches(0.04),
                 status, size=9, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        # Priority pill
        prio_left = left + col_widths[0] + col_widths[1] + col_widths[2] + Inches(0.1)
        prio_w = col_widths[3] - Inches(0.2)
        prio_fill = {
            "Highest": RED_FILL, "High": YELLOW_FILL, "Med": HEADER_FILL,
            "Low": HEADER_FILL, "Maintain": GREEN_FILL, "—": WHITE,
        }.get(prio, HEADER_FILL)
        prio_text_color = {
            "Highest": RED_DOT, "High": RGBColor(0x8B, 0x5A, 0x00),
        }.get(prio, INK)
        add_rect(slide, prio_left, row_top + Inches(0.07), prio_w,
                 row_h - Inches(0.14), prio_fill, line_color=RULE)
        add_text(slide, prio_left, row_top + Inches(0.05), prio_w,
                 row_h - Inches(0.1), prio, size=9, bold=True,
                 color=prio_text_color, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)


def build_slide_4_scorecard(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "Our scorecard", eyebrow="vs Anthropic's 4-layer stack",
                 page_num=page, total=total)
    rows = [
        ("y", "Canonical datasets",
         "etoro_kpi views own MIMO / AUM / PFOF; other domains overlap with DDR + Synapse TVFs",
         "Med"),
        ("y", "Skills + models colocation",
         "Both live in DataPlatform; missing CI hook enforcing skill-touch on model PRs",
         "High"),
        ("g", "UC metadata as product",
         "Column-level descriptions deployed across 6 domains; ambivalent on retrieval lift",
         "Low"),
        ("y", "Semantic layer (declarative)",
         "No declarative metric layer — but skills route to etoro_kpi canonical views first",
         "High"),
        ("g", "Lineage + table ranking",
         "Genie Code has UC lineage access today",
         "—"),
        ("y", "Query corpus",
         "Captured (MCP + Genie gateway); distillation manual today",
         "Med"),
        ("y", "Business context",
         "SME / TVF docs in Synapse Wiki + Confluence; not piped to agents",
         "Low"),
        ("g", "Skills (knowledge router)",
         "Hub-and-spoke; ~13 entry + 45 sub-skills; MCP-served, CI-validated",
         "Maintain"),
        ("y", "Skills (unbook / process)",
         "Substance authored (3-skill triple) — analyst-triggered by design today, not auto-fired on every question. Open H2 decision: promote a minimal subset?",
         "Open"),
        ("r", "Offline evals",
         "None today — but /feedback app captures graded Q&A in production",
         "Highest"),
        ("r", "Ablation methodology",
         "None — gated on evals",
         "Med"),
        ("r", "Provenance footer",
         "No footer on Genie / MCP responses",
         "High"),
        ("r", "Adversarial review on every answer",
         "Deliberately NOT adopted — +32% tokens / +72% latency / +6% accuracy doesn't pencil. Pattern kept in library for manual high-stakes use.",
         "Reject"),
        ("g", "Passive monitoring",
         "genie_audit_events + MCP gateway logs live",
         "Med"),
        ("y", "Active correction harvesting",
         "Substrate exists (/feedback + MCP user-message logs); missing scheduled classifier + PR-draft agent",
         "High"),
    ]
    add_scorecard_table(s, rows, Inches(0.5), Inches(1.4),
                        Inches(12.3), Inches(5.5))
    # Score legend
    add_text(s, Inches(0.5), Inches(6.95), Inches(12), Inches(0.4),
             "Score: 5 green   ·   6 yellow   ·   4 red       "
             "1 red is a DELIBERATE non-adoption with rationale, not a gap (see Slide 6).",
             size=11, italic=True, color=SUBTLE)


def build_slide_5_right(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "What we're doing RIGHT",
                 eyebrow="Where we measure up — or surpass",
                 page_num=page, total=total)
    items = [
        ("Knowledge skill corpus. ",
         "Hub-and-spoke routing, CI-enforced frontmatter, kebab-case names, required body sections. Identical shape to Anthropic's appendix skeleton."),
        ("Unbook SUBSTANCE in production — decomposed even better than Anthropic's. ",
         "Three-skill triple on Databricks Assistant (data-analysis-playbook + data-analysis-patterns + data-analysis-pattern-library) splits process from routing from detail. Caveat: analyst-triggered today by design. Anthropic's lift comes from auto-firing; that decision is open for us (Slide 7)."),
        ("Skills + models in the same repo (DataPlatform). ",
         "CI-deployed. Same colocation principle Anthropic endorses — what's still missing is the enforcement hook."),
        ("Cross-surface portability. ",
         "Same skill served via MCP gateway → Cursor IDE, Genie Code, standalone agents."),
        ("Telemetry foundations live. ",
         "genie_audit_events + monitoring_mcp_logs_mcp_gateway capture skill loads, NL prompts, generated SQL, query-history joins."),
        ("/feedback Databricks app — strategic asset. ",
         "Every Genie answer one-click graded; landed in de_output_genie_code_skill_feedback. Fastest path to a labeled eval set in the industry — Anthropic builds evals by hand; we harvest them in production with grades."),
        ("UC as a documented warehouse. ",
         "~10k+ column comments deployed across 6 domains."),
        ("Workspace-level assistant defaults exist. ",
         ".assistant_workspace_instructions.md at workspace root is the lever for cross-Genie-Code behavior — already where the analysis triple is anchored."),
    ]
    add_bullet_list(s, Inches(0.6), Inches(1.35), Inches(12.2), Inches(5.5),
                    items, size=13, spacing=10)


def build_slide_6_gaps(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "What we're doing WRONG", eyebrow="The four live gaps",
                 page_num=page, total=total)
    gaps = [
        ("Gap 1 — We cannot measure our accuracy yet",
         "Zero pinned Q&A evals. We don't know if our agents are at 21% or 95%. Every gap below is invisible without this. But: /feedback already captures graded Q&A — path is to synthesize the eval set from that table, not author from scratch."),
        ("Gap 2 — Skill ↔ model drift has no CI guard",
         "Anthropic: 90% of their data-model PRs include a skill-file change in the same diff, enforced by CI. We have colocation but no enforcement. Every canonical-view change silently rots the skill that points at it."),
        ("Gap 3 — No provenance, no online safety net",
         "Every Anthropic response carries a footer: source tier / freshness / owner / skill used. We ship raw numbers with no signal. Silent wrong answers are our highest-risk failure mode and we have zero mitigation. KEY INSIGHT: our existing final-answer-assembly pattern IS the provenance footer — Gap 3 collapses into one cheap fix (Slide 7+9)."),
        ("Gap 4 — No active correction loop (yet)",
         "Anthropic: scheduled agent scans channels every few hours for correction language → drafts skill-file PRs → tags domain owner. We have the substrate (/feedback + MCP user-message logs); MCP harvest is ready to start today; Genie Code via our enrichment; classic Genie waits for Databricks logs (out of our control)."),
    ]
    y = Inches(1.4)
    for i, (lead, body) in enumerate(gaps):
        h = Inches(1.30)
        add_rect(s, Inches(0.6), y, Inches(0.55), Inches(0.55), RED_DOT)
        add_text(s, Inches(0.6), y + Inches(0.02), Inches(0.55), Inches(0.55),
                 str(i + 1), size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.3), y, Inches(11.5), Inches(0.32), lead,
                 size=14, bold=True, color=NAVY)
        add_text(s, Inches(1.3), y + Inches(0.34), Inches(11.5),
                 h - Inches(0.3), body, size=11, color=INK)
        y += h + Inches(0.08)


def build_slide_deliberate_non_adoption(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "What we're DELIBERATELY NOT adopting",
                 eyebrow="One Anthropic recommendation, rejected with rationale",
                 page_num=page, total=total)

    # Big header box
    add_rect(s, Inches(0.6), Inches(1.3), Inches(12.2), Inches(0.7),
             RED_FILL, line_color=RED_DOT, line_width=Pt(1.5))
    add_text(s, Inches(0.85), Inches(1.40), Inches(11.6), Inches(0.5),
             "Mandatory adversarial review on every user-facing answer",
             size=18, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.6), Inches(2.15), Inches(12), Inches(0.4),
             "Anthropic enforces a Challenge-the-Solution sub-agent call on every analytical answer:",
             size=12, color=INK)
    # Cost / benefit table
    metrics = [
        ("Accuracy", "+6%", GREEN_DOT),
        ("Tokens", "+32%", YELLOW_DOT),
        ("Latency", "+72%", RED_DOT),
    ]
    for i, (k, v, col) in enumerate(metrics):
        left = Inches(0.6 + i * 4.1)
        add_rect(s, left, Inches(2.6), Inches(3.9), Inches(0.9), BAND_FILL,
                 line_color=RULE)
        add_text(s, left + Inches(0.2), Inches(2.65), Inches(2.0),
                 Inches(0.35), k, size=11, color=SUBTLE)
        add_text(s, left + Inches(0.2), Inches(2.92), Inches(3.5),
                 Inches(0.5), v, size=22, bold=True, color=col)

    # Position
    add_text(s, Inches(0.6), Inches(3.7), Inches(12), Inches(0.4),
             "Our position", size=13, bold=True, color=ACCENT)
    add_text(s, Inches(0.6), Inches(4.05), Inches(12.2), Inches(1.0),
             "The economics don't pencil at current model costs. A nearly-doubled response time on every "
             "Genie / MCP query — for a 6-point accuracy lift that's likely smaller for our use case because "
             "most of our questions are KPI-style, not open-ended diagnostics — is not a trade users will "
             "accept on day-to-day questions.",
             size=12, color=INK)

    add_text(s, Inches(0.6), Inches(5.05), Inches(12), Inches(0.4),
             "What we keep", size=13, bold=True, color=ACCENT)
    add_text(s, Inches(0.6), Inches(5.40), Inches(12.2), Inches(0.7),
             "The Challenge-the-Solution pattern stays in data-analysis-pattern-library and is invoked MANUALLY "
             "by analysts on high-stakes analyses (board metrics, regulator-facing numbers, model-validation work) — "
             "same procedure, applied where the latency cost is justified.",
             size=12, color=INK)

    add_text(s, Inches(0.6), Inches(6.20), Inches(12), Inches(0.4),
             "Revisit", size=13, bold=True, color=ACCENT)
    add_text(s, Inches(0.6), Inches(6.55), Inches(12.2), Inches(0.5),
             "Every 6 months as model cost / latency curves move. If Claude latency drops 3x or eval data shows "
             "our domain accuracy is below 80%, this becomes a \"yes.\"",
             size=12, italic=True, color=INK)


def build_slide_open_decision(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "The open H2 design decision",
                 eyebrow="Should we promote a minimal unbook subset to auto-fire?",
                 page_num=page, total=total)

    # Pull-quote
    add_rect(s, Inches(0.6), Inches(1.3), Inches(12.2), Inches(0.7),
             HEADER_FILL)
    add_text(s, Inches(0.85), Inches(1.40), Inches(11.7), Inches(0.55),
             "Anthropic's unbook fires on every question. Ours is analyst-triggered. "
             "Should we promote a minimal subset to auto-fire?",
             size=14, italic=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.6), Inches(2.15), Inches(12), Inches(0.4),
             "Three candidates, ordered by cost",
             size=13, bold=True, color=NAVY)

    # Table
    headers = ["Candidate", "What fires on every question", "Token cost", "Latency cost", "What it buys"]
    col_widths = [Inches(2.0), Inches(4.5), Inches(1.3), Inches(1.5), Inches(2.9)]
    top = Inches(2.55)
    header_h = Inches(0.40)
    add_rect(s, Inches(0.6), top, sum(col_widths, Inches(0)), header_h, NAVY)
    x = Inches(0.6)
    for w, h in zip(col_widths, headers):
        add_text(s, x + Inches(0.08), top + Inches(0.05),
                 w - Inches(0.16), header_h - Inches(0.1), h,
                 size=10, bold=True, color=WHITE)
        x += w

    rows = [
        ("(a) final-answer-assembly",
         "Mandatory output template: lead with answer + confidence H/M/L + metric definition + provenance footer.",
         "+5–10%", "Negligible",
         "Subsumes Gap 3. Cheapest. RECOMMENDED START."),
        ("(b) + metric-definition-check",
         "At deliver-time: re-state the metric the query computes; flag if it diverges from likely user intent.",
         "+10–15%", "+5–10s",
         "Attacks ambiguity at output time, not just at planning."),
        ("(c) + question-framing",
         "At intake: surface ambiguous terms (\"last week\"? \"active\"?) — ask one targeted clarifying question.",
         "+15–25%", "+10–20s on ambiguous Qs",
         "Closest to Anthropic's pattern — clarify before answering."),
    ]
    row_top = top + header_h
    row_h = Inches(0.95)
    for i, row in enumerate(rows):
        rt = row_top + i * row_h
        fill = BAND_FILL if i % 2 == 0 else WHITE
        # Highlight (a)
        if i == 0:
            fill = GREEN_FILL
        add_rect(s, Inches(0.6), rt, sum(col_widths, Inches(0)), row_h, fill,
                 line_color=RULE)
        x = Inches(0.6)
        for j, (w, val) in enumerate(zip(col_widths, row)):
            bold = (j == 0)
            add_text(s, x + Inches(0.08), rt + Inches(0.05),
                     w - Inches(0.16), row_h - Inches(0.1), val,
                     size=10, bold=bold, color=INK,
                     anchor=MSO_ANCHOR.MIDDLE)
            x += w

    # Recommendation + why
    rec_top = Inches(5.95)
    add_rect(s, Inches(0.6), rec_top, Inches(12.2), Inches(1.0),
             HEADER_FILL, line_color=RULE)
    add_text(s, Inches(0.85), rec_top + Inches(0.08), Inches(11.7),
             Inches(0.35), "Recommendation",
             size=12, bold=True, color=ACCENT)
    add_text(s, Inches(0.85), rec_top + Inches(0.40), Inches(11.7),
             Inches(0.6),
             "Start with (a). It's a routing change, not content authoring — data-analysis-patterns already has a "
             "\"Communication-only\" output mode. Measure via the eval gate (Slide 9). Promote to (b) or (c) only if data justifies it. "
             "The 7-step heavyweight loop stays analyst-triggered for non-trivial analyses.",
             size=11, color=INK)


def build_slide_7_ambivalent_and_designed(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "Ambivalent investments + what we deferred by design",
                 eyebrow="Don't sweat these; revisit one",
                 page_num=page, total=total)

    # Ambivalent
    add_rect(s, Inches(0.6), Inches(1.3), Inches(6.0), Inches(5.4), BAND_FILL,
             line_color=RULE)
    add_text(s, Inches(0.85), Inches(1.4), Inches(5.5), Inches(0.4),
             "AMBIVALENT (DON'T SWEAT)", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(0.85), Inches(1.85), Inches(5.5), Inches(0.45),
             "UC column descriptions",
             size=16, bold=True, color=NAVY)
    add_text(s, Inches(0.85), Inches(2.35), Inches(5.5), Inches(4.0),
             "Heavy investment to date. Anthropic explicitly says docs help "
             "but aren't the bottleneck once a semantic layer + skill router exist — "
             "they confirmed via ablation that giving the agent grep over "
             "thousands of analyst notebooks moved accuracy <1pt.\n\n"
             "Verdict: keep what we have, stop investing net-new effort "
             "here until evals show the gap. Reallocate the column-doc "
             "team-time to evals + correction harvesting.",
             size=12, color=INK)

    # Deferred by design
    add_rect(s, Inches(6.85), Inches(1.3), Inches(6.0), Inches(5.4), HEADER_FILL,
             line_color=RULE)
    add_text(s, Inches(7.10), Inches(1.4), Inches(5.5), Inches(0.4),
             "DEFERRED BY DESIGN (REVISIT)", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(7.10), Inches(1.85), Inches(5.5), Inches(0.45),
             "Declarative semantic layer (UC metric views)",
             size=16, bold=True, color=NAVY)
    add_text(s, Inches(7.10), Inches(2.35), Inches(5.5), Inches(4.2),
             "We have etoro_kpi.* canonical prep views + skill-based routing. "
             "For ~5 headline KPIs the agent already lands on a canonical "
             "view without a declarative metric primitive.\n\n"
             "Missing: the compilation step. In a UC-metric-views world the "
             "agent calls metric.measure(\"revenue\", group_by=\"region\") "
             "and gets one number — same number every BI surface produces.\n\n"
             "Anthropic's data: this is their highest-trust path and the "
             "mandatory default. Economics changed when UC metric views "
             "went GA in 2026. A 10-KPI pilot would compress ambiguity on "
             "the most-asked questions.",
             size=12, color=INK)


def build_slide_8_evals(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "The shortcut: synthesize evals from the /feedback app",
                 eyebrow="Our single biggest H2 win",
                 page_num=page, total=total)

    add_rect(s, Inches(0.6), Inches(1.3), Inches(12.2), Inches(0.7),
             HEADER_FILL)
    add_text(s, Inches(0.85), Inches(1.4), Inches(11.7), Inches(0.55),
             "\"Pin every eval to a snapshot date, write it against a stable fact table, "
             "or have the grader judge the agent's query rather than its number.\"  — Anthropic",
             size=12, italic=True, color=INK)

    add_text(s, Inches(0.6), Inches(2.15), Inches(12), Inches(0.4),
             "/feedback app already captures, per Genie submission:",
             size=13, bold=True, color=NAVY)
    feedback_fields = [
        "NL question",
        "Skills the agent loaded",
        "SQL the agent generated",
        "Numeric result",
        "User grade (1–5 ★)",
        "Optional free-text correction",
    ]
    for i, f in enumerate(feedback_fields):
        col = i % 3
        row = i // 3
        left = Inches(0.7 + col * 4.1)
        top = Inches(2.6 + row * 0.4)
        add_rect(s, left, top, Inches(0.18), Inches(0.18), ACCENT)
        add_text(s, left + Inches(0.3), top - Inches(0.02), Inches(3.7),
                 Inches(0.4), f, size=12, color=INK)

    add_text(s, Inches(0.6), Inches(3.55), Inches(12), Inches(0.4),
             "Pareto eval-set construction:",
             size=13, bold=True, color=NAVY)
    steps = [
        ("4★ + 5★ submissions → graded-correct canonical Q&A pairs.",
         " Pin to created_at snapshot. Grader checks SQL shape, not number."),
        ("1★ + 2★ submissions with corrections → correction-harvest source.",
         " Each is a candidate skill-file edit."),
        ("Frequency-weight by NL similarity (LLM clustering).",
         " Top 30 question clusters per domain cover the long tail."),
        ("Land eval runs as telemetry in a Delta table.",
         " skill_version_sha, model_id, run_ts, passed_bool, per_assertion_json, tokens, latency."),
        ("Wire into CI on every skill PR.",
         " Run only the eval slice affected by the diff."),
        ("Gate at 90% per domain (Anthropic's threshold).",
         " No stakeholder announcement below the gate."),
    ]
    y = Inches(4.0)
    for i, (lead, body) in enumerate(steps):
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.36),
                                   Inches(0.36))
        circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
        circ.line.fill.background(); circ.shadow.inherit = False
        add_text(s, Inches(0.6), y, Inches(0.36), Inches(0.36), str(i + 1),
                 size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rich(s, Inches(1.10), y + Inches(0.04), Inches(11.5), Inches(0.36), [
            {"runs": [
                {"text": lead, "size": 11, "bold": True, "color": INK},
                {"text": body, "size": 11, "color": INK},
            ]},
        ])
        y += Inches(0.42)


def build_slide_9_provenance(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "Provenance footer — surfaces in scope",
                 eyebrow="Cheapest mitigation for silent wrong answers",
                 page_num=page, total=total)

    # Footer template box
    add_rect(s, Inches(0.6), Inches(1.3), Inches(12.2), Inches(1.55),
             RGBColor(0x14, 0x1F, 0x33))
    add_text(s, Inches(0.85), Inches(1.4), Inches(11.7), Inches(0.32),
             "Footer template Anthropic appends to every answer:",
             size=11, color=RGBColor(0xB7, 0xCF, 0xE6))
    add_text(s, Inches(0.85), Inches(1.72), Inches(11.7), Inches(1.1),
             "Source: semantic layer / curated view / raw exploration\n"
             "Freshness: data through YYYY-MM-DD\n"
             "Owner: <team>\n"
             "Skill used: <skill_id>@<commit_sha>",
             size=12, color=WHITE)

    add_text(s, Inches(0.6), Inches(3.0), Inches(12), Inches(0.4),
             "Implementation paths — only the surfaces we control",
             size=13, bold=True, color=NAVY)

    headers = ["Surface", "How", "Status"]
    col_widths = [Inches(3.7), Inches(6.6), Inches(2.0)]
    rows = [
        ("MCP gateway responses",
         "Append footer in the gateway middleware. We own the layer end-to-end.",
         "Easy · 1 PR"),
        ("Genie Code (Databricks Assistant in notebooks)",
         "Two options: (a) add the mandate to .assistant_workspace_instructions.md "
         "at workspace root — applies globally; or (b) push a dedicated "
         "provenance-footer skill into /.assistant/skills/.",
         "Medium"),
        ("Cursor IDE",
         "SKIP — only our custom MCP sees skill-load context anyway.",
         "n/a — covered by MCP"),
        ("Classic Genie Space (chat surface)",
         "SKIP — no skill / instruction injection point yet. Databricks-controlled. "
         "Wait for product.",
         "Out of scope"),
    ]
    # header
    x = Inches(0.6)
    top = Inches(3.45)
    header_h = Inches(0.36)
    add_rect(s, x, top, sum(col_widths, Inches(0)), header_h, NAVY)
    for w, h in zip(col_widths, headers):
        add_text(s, x + Inches(0.1), top + Inches(0.05), w - Inches(0.2),
                 header_h - Inches(0.1), h, size=11, bold=True, color=WHITE)
        x += w
    # rows
    for i, (surface, how, status) in enumerate(rows):
        row_top = top + header_h + Inches(0.55 * i)
        row_h = Inches(0.85)
        skipped = "SKIP" in how
        fill = BAND_FILL if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.6), row_top, sum(col_widths, Inches(0)), row_h,
                 fill, line_color=RULE)
        # surface
        add_text(s, Inches(0.7), row_top + Inches(0.05),
                 col_widths[0] - Inches(0.1), row_h - Inches(0.1),
                 surface, size=11, bold=True,
                 color=SUBTLE if skipped else INK,
                 anchor=MSO_ANCHOR.MIDDLE,
                 italic=skipped)
        add_text(s, Inches(0.7) + col_widths[0], row_top + Inches(0.05),
                 col_widths[1] - Inches(0.2), row_h - Inches(0.1),
                 how, size=10, color=SUBTLE if skipped else INK,
                 anchor=MSO_ANCHOR.MIDDLE,
                 italic=skipped)
        # status pill
        pill_left = Inches(0.7) + col_widths[0] + col_widths[1]
        pill_w = col_widths[2] - Inches(0.2)
        pill_fill = HEADER_FILL if skipped else GREEN_FILL if "Easy" in status else YELLOW_FILL
        add_rect(s, pill_left, row_top + Inches(0.22), pill_w,
                 row_h - Inches(0.44), pill_fill, line_color=RULE)
        add_text(s, pill_left, row_top + Inches(0.2), pill_w,
                 row_h - Inches(0.4), status, size=10, bold=True,
                 color=SUBTLE if skipped else INK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def build_slide_machine(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "The machine — from static corpus to living business brain",
                 eyebrow="The big idea",
                 page_num=page, total=total)

    # Pull-quote
    add_rect(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.6),
             HEADER_FILL)
    add_text(s, Inches(0.7), Inches(1.32), Inches(11.9), Inches(0.5),
             "The point of H2 is not to ship more skills. It's to stop authoring them as one-off artifacts "
             "and start running them as the output of an autonomous system.",
             size=12, italic=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # Three columns: Inputs / Orchestration / Outputs
    diagram_top = Inches(2.0)
    diagram_h = Inches(3.85)
    col_w = Inches(3.95)
    gap = Inches(0.10)
    arrow_w = Inches(0.30)

    # Column 1 — Inputs (sensors)
    in_left = Inches(0.50)
    add_rect(s, in_left, diagram_top, col_w, Inches(0.40), ACCENT)
    add_text(s, in_left, diagram_top, col_w, Inches(0.40),
             "INPUTS — sensors", size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, in_left, diagram_top + Inches(0.40), col_w,
             diagram_h - Inches(0.40), BAND_FILL, line_color=RULE)
    inputs = [
        "New Confluence docs",
        "New SharePoint docs",
        "UC schema changes (information_schema + UC lineage)",
        "DataPlatform PRs touching canonical models",
        "MCP query telemetry",
        "Genie audit events",
        "/feedback 1–2★ grades + correction comments",
    ]
    y = diagram_top + Inches(0.55)
    for item in inputs:
        # dot
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, in_left + Inches(0.20),
                                  y + Inches(0.05), Inches(0.10), Inches(0.10))
        circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
        circ.line.fill.background(); circ.shadow.inherit = False
        add_text(s, in_left + Inches(0.40), y - Inches(0.02),
                 col_w - Inches(0.50), Inches(0.45), item, size=10,
                 color=INK)
        y += Inches(0.45)

    # Arrow 1
    arrow_y = diagram_top + Inches(1.6)
    arr1 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               in_left + col_w + gap,
                               arrow_y, arrow_w, Inches(0.55))
    arr1.fill.solid(); arr1.fill.fore_color.rgb = NAVY
    arr1.line.fill.background(); arr1.shadow.inherit = False

    # Column 2 — Orchestration
    orch_left = in_left + col_w + gap + arrow_w + gap
    add_rect(s, orch_left, diagram_top, col_w, Inches(0.40), NAVY)
    add_text(s, orch_left, diagram_top, col_w, Inches(0.40),
             "ORCHESTRATION — judges + routers", size=11, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, orch_left, diagram_top + Inches(0.40), col_w,
             diagram_h - Inches(0.40), RGBColor(0xF0, 0xF4, 0xFB),
             line_color=RULE)
    orch_items = [
        ("Scheduled LLM classifiers + routers",
         "Detect skill drift, correction language, schema deltas"),
        ("Eval gate (90% per domain)",
         "Per-domain accuracy SLA — refuses regressions"),
        ("Skill-touch CI hook",
         "Model PR must touch matching skill file"),
        ("Frequency-weighted Q-cluster prioritizer",
         "Focuses effort on the long-tail questions"),
    ]
    y = diagram_top + Inches(0.55)
    for h, body in orch_items:
        add_text(s, orch_left + Inches(0.15), y, col_w - Inches(0.30),
                 Inches(0.30), h, size=10, bold=True, color=NAVY)
        add_text(s, orch_left + Inches(0.15), y + Inches(0.25),
                 col_w - Inches(0.30), Inches(0.45), body, size=9,
                 color=SUBTLE)
        y += Inches(0.78)

    # Arrow 2
    arr2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               orch_left + col_w + gap,
                               arrow_y, arrow_w, Inches(0.55))
    arr2.fill.solid(); arr2.fill.fore_color.rgb = NAVY
    arr2.line.fill.background(); arr2.shadow.inherit = False

    # Column 3 — Outputs
    out_left = orch_left + col_w + gap + arrow_w + gap
    add_rect(s, out_left, diagram_top, col_w, Inches(0.40), GREEN_DOT)
    add_text(s, out_left, diagram_top, col_w, Inches(0.40),
             "OUTPUTS — actions", size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, out_left, diagram_top + Inches(0.40), col_w,
             diagram_h - Inches(0.40), BAND_FILL, line_color=RULE)
    outputs = [
        "Draft skill PRs (auto-tagged to domain owner)",
        "Eval-set additions (from /feedback)",
        "Schema-drift alerts (to domain owner)",
        "Stale-skill flags (source doc deprecated / table retired)",
        "Accuracy dashboard per domain",
    ]
    y = diagram_top + Inches(0.65)
    for item in outputs:
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, out_left + Inches(0.20),
                                  y + Inches(0.05), Inches(0.10),
                                  Inches(0.10))
        circ.fill.solid(); circ.fill.fore_color.rgb = GREEN_DOT
        circ.line.fill.background(); circ.shadow.inherit = False
        add_text(s, out_left + Inches(0.40), y - Inches(0.02),
                 col_w - Inches(0.50), Inches(0.55), item, size=10,
                 color=INK)
        y += Inches(0.55)

    # Three rules of the machine, bottom band
    rules_top = Inches(6.05)
    add_rect(s, Inches(0.5), rules_top, Inches(12.3), Inches(0.85),
             RGBColor(0x14, 0x1F, 0x33))
    add_text(s, Inches(0.7), rules_top + Inches(0.07), Inches(11.9),
             Inches(0.30), "Three rules of the machine",
             size=11, bold=True, color=RGBColor(0x6F, 0xB0, 0xD9))
    rules = [
        ("Every input becomes a trigger.",
         " A new Confluence page is a signal to verify the routing."),
        ("Every trigger becomes a PR or eval.",
         " Humans review drafts; the machine does the boring work."),
        ("The eval gate is the brain.",
         " Nothing ships unless the domain's slice clears 90%."),
    ]
    for i, (lead, body) in enumerate(rules):
        rl = Inches(0.7 + i * 4.1)
        rt = rules_top + Inches(0.38)
        add_rich(s, rl, rt, Inches(4.0), Inches(0.50), [
            {"runs": [
                {"text": f"{i+1}. ", "size": 10, "bold": True, "color": WHITE},
                {"text": lead, "size": 10, "bold": True, "color": WHITE},
                {"text": body, "size": 10, "color": RGBColor(0xB7, 0xCF, 0xE6)},
            ]},
        ])


def build_slide_10_h2_plan(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "H2 plan — building the machine, component by component",
                 eyebrow="Each priority is one part of the system on the prior slide",
                 page_num=page, total=total)

    # Components-of-the-machine table
    headers = ["#", "Component", "What it does", "Status today"]
    col_widths = [Inches(0.55), Inches(2.7), Inches(6.6), Inches(2.4)]
    top = Inches(1.35)
    header_h = Inches(0.40)
    total_w = sum(col_widths, Inches(0))

    add_rect(s, Inches(0.5), top, total_w, header_h, NAVY)
    x = Inches(0.5)
    for w, h in zip(col_widths, headers):
        align = PP_ALIGN.CENTER if h == "#" else PP_ALIGN.LEFT
        add_text(s, x + Inches(0.08), top + Inches(0.05),
                 w - Inches(0.16), header_h - Inches(0.1), h,
                 size=11, bold=True, color=WHITE, align=align)
        x += w

    rows = [
        ("1", ACCENT,
         "TRUTH SENSOR",
         "Eval substrate",
         "Synthesize eval set from /feedback (4★+5★ → canonical pairs; 1★+2★ → corrections). CI gate at 90% per domain.",
         "Substrate live in /feedback app; harvester + gate logic to be built."),
        ("2", ACCENT,
         "MODEL-CHANGE WATCHER",
         "Skill-touch CI hook",
         "Block any DataPlatform PR that changes canonical tables / prep views without touching the matching skill file.",
         "New, single PR."),
        ("3", ACCENT,
         "OUTPUT CONTRACT",
         "final-answer-assembly enforced",
         "Every answer ships with confidence + metric definition + source tier + skill_id@sha (MCP + Genie Code via assistant instructions).",
         "Pattern authored in library; routing change only."),
        ("4", RED_DOT,
         "MULTI-SOURCE WATCHER FLEET",
         "the heart of the machine",
         "Scheduled jobs that watch and trigger:  (a) MCP user-message logs + /feedback 1–2★ for corrections,  "
         "(b) new Confluence pages tagged to a domain,  (c) new SharePoint docs in mapped folders,  "
         "(d) UC schema changes via information_schema + UC lineage deltas.  Each → LLM classifier → draft skill-file PR tagged to domain owner.",
         "NEW. MCP-correction strand can start now (logs are live); Confluence + SharePoint + UC strands reuse the same PR-draft framework."),
        ("5", ACCENT,
         "ABLATION-GRADE TELEMETRY",
         "(bonus)",
         "Enrich monitoring_mcp_logs_mcp_gateway + Genie sibling table with skill_version_sha, model_id, token_in/out, latency_ms.",
         "New columns + writer changes."),
        ("?", YELLOW_DOT,
         "OPEN DECISION",
         "10-KPI UC-metric-view pilot",
         "Yes / no on a thin declarative semantic layer on top of etoro_kpi. Lift over canonical-view routing isn't proven for our shape.",
         "Decision in H2 planning."),
    ]
    row_top = top + header_h
    row_h_default = Inches(0.78)
    row_h_big = Inches(1.10)
    for i, row in enumerate(rows):
        n, dot_col, label, subtitle, what, status = row
        is_big = (i == 3)
        row_h = row_h_big if is_big else row_h_default
        fill = BAND_FILL if i % 2 == 0 else WHITE
        if is_big:
            fill = RGBColor(0xFF, 0xEC, 0xEF)
        add_rect(s, Inches(0.5), row_top, total_w, row_h, fill,
                 line_color=RULE)
        # Number circle
        circ_x = Inches(0.5) + Inches(0.10)
        circ_y = row_top + (row_h - Inches(0.40)) / 2
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, circ_x, circ_y,
                                  Inches(0.40), Inches(0.40))
        circ.fill.solid(); circ.fill.fore_color.rgb = dot_col
        circ.line.fill.background(); circ.shadow.inherit = False
        add_text(s, circ_x, circ_y, Inches(0.40), Inches(0.40), n,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Component column — label + subtitle stacked
        comp_x = Inches(0.5) + col_widths[0] + Inches(0.05)
        add_text(s, comp_x, row_top + Inches(0.06),
                 col_widths[1] - Inches(0.1), Inches(0.30), label,
                 size=9, bold=True, color=dot_col)
        add_text(s, comp_x, row_top + Inches(0.30),
                 col_widths[1] - Inches(0.1), row_h - Inches(0.32),
                 subtitle, size=11, bold=True, color=INK)
        # What it does
        what_x = comp_x + col_widths[1]
        add_text(s, what_x, row_top + Inches(0.06),
                 col_widths[2] - Inches(0.10), row_h - Inches(0.12),
                 what, size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        # Status
        stat_x = what_x + col_widths[2]
        add_text(s, stat_x, row_top + Inches(0.06),
                 col_widths[3] - Inches(0.10), row_h - Inches(0.12),
                 status, size=9, italic=True, color=SUBTLE,
                 anchor=MSO_ANCHOR.MIDDLE)
        row_top += row_h

    # Not-on-the-list footer
    add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.30),
             "NOT on the list: adversarial review enforcement (Slide 6 — deliberate non-adoption); "
             "promoting the full unbook triple to auto-fire (Slide 7 — open beyond candidate (a)).",
             size=10, italic=True, color=SUBTLE)


def build_slide_11_asks(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "Asks of leadership", eyebrow="What we need from you",
                 page_num=page, total=total)

    # Primary ask — boxed
    add_rect(s, Inches(0.6), Inches(1.3), Inches(12.2), Inches(2.6),
             RED_FILL, line_color=RED_DOT, line_width=Pt(1.5))
    add_text(s, Inches(0.85), Inches(1.4), Inches(11.7), Inches(0.4),
             "PRIMARY ASK", size=11, bold=True, color=RED_DOT)
    add_text(s, Inches(0.85), Inches(1.78), Inches(11.7), Inches(0.5),
             "Assign domain owners and commit them to the eval gate",
             size=18, bold=True, color=NAVY)
    add_text(s, Inches(0.85), Inches(2.35), Inches(11.7), Inches(0.4),
             "The machine has one human-in-the-loop role that cannot be automated: the domain owner who "
             "adjudicates ground truth for their slice. Per Anthropic — non-negotiable.",
             size=11, italic=True, color=INK)
    add_rich(s, Inches(0.85), Inches(2.85), Inches(11.7), Inches(1.0), [
        {"runs": [
            {"text": "• Internal teams ", "size": 11, "bold": True, "color": INK},
            {"text": "(Data Analytics, Product Analytics — inside our org): ",
             "size": 11, "color": INK},
            {"text": "can be mandated. ", "size": 11, "bold": True, "color": RED_DOT},
            {"text": "Domain ownership assigned; team signs off on the eval slice.",
             "size": 11, "color": INK},
        ]},
        {"runs": [
            {"text": "• Cross-functional business domains ",
             "size": 11, "bold": True, "color": INK},
            {"text": "(Marketing, Trading Ops, Payments, US Ops, Compliance…): ",
             "size": 11, "color": INK},
            {"text": "needs leadership-level comms. ",
             "size": 11, "bold": True, "color": RED_DOT},
            {"text": "Message: \"the agent will give your users wrong answers unless your team owns the eval slice. "
             "We build the substrate; you adjudicate ground truth.\" Without this we stall at ~75%.",
             "size": 11, "color": INK},
        ], "space_before": 4},
    ])

    # Other asks
    add_text(s, Inches(0.6), Inches(4.10), Inches(12), Inches(0.4),
             "And three smaller asks", size=13, bold=True, color=NAVY)
    other_asks = [
        ("2.", "Greenlight the machine's first watchers",
         "Skill-touch CI hook on DataPlatform (one PR, no product impact) + MCP correction harvester (scheduled job on logs that already exist — ready to start now)."),
        ("3.", "Commit a baseline accuracy number for H2",
         "Suggested floor: 75% on top 4 domains by Q3, 90% by Q4. The machine without a target is just plumbing."),
        ("4.", "Decision on the 10-KPI UC-metric-view pilot",
         "Yes / no on a thin semantic-layer pilot in H2. We have the canonical views; this is the declarative wrapper on top."),
    ]
    y = Inches(4.55)
    for n, lead, body in other_asks:
        add_text(s, Inches(0.7), y, Inches(0.5), Inches(0.4), n,
                 size=15, bold=True, color=ACCENT)
        add_text(s, Inches(1.2), y, Inches(11.5), Inches(0.35), lead,
                 size=12, bold=True, color=NAVY)
        add_text(s, Inches(1.2), y + Inches(0.32), Inches(11.5),
                 Inches(0.5), body, size=11, color=INK)
        y += Inches(0.78)


def build_slide_12_summary(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "One-page summary", eyebrow="Take this back to the room",
                 page_num=page, total=total)

    sections = [
        ("The big idea",
         "We're done authoring skills as static artifacts. H2 is about building THE MACHINE — an autonomous knowledge system that watches Confluence, SharePoint, UC schema, MCP/Genie telemetry, and /feedback grades, and turns every signal into a draft skill PR or an eval-set entry, gated by a per-domain accuracy SLA."),
        ("Where we are",
         "5 green / 6 yellow / 4 red against Anthropic's framework — 1 red is a deliberate non-adoption. Skills + telemetry + repo colocation + the /feedback app are the real strengths."),
        ("Two assets most teams don't have",
         "(1) /feedback app — eval-building, solved.   (2) Three-skill analysis triple — substance of an unbook, decomposed better than Anthropic's. Analyst-triggered today; we may promote one cheap subset to auto-fire."),
        ("One thing we're rejecting on purpose",
         "Mandatory adversarial review on every answer — economics don't pencil (+72% latency for +6% accuracy). Pattern stays for manual high-stakes use."),
        ("Five components of the machine (Slide 13)",
         "(1) Truth sensor — eval substrate · (2) Model-change watcher · (3) Output contract · (4) MULTI-SOURCE WATCHER FLEET (MCP corrections + Confluence + SharePoint + UC schema deltas) — the heart of the machine · (5) Ablation-grade telemetry. Open: 10-KPI metric-view pilot."),
    ]
    y = Inches(1.30)
    for h, body in sections:
        add_rect(s, Inches(0.6), y, Inches(12.2), Inches(1.00), BAND_FILL,
                 line_color=RULE)
        add_text(s, Inches(0.85), y + Inches(0.08), Inches(11.7),
                 Inches(0.30), h, size=13, bold=True, color=ACCENT)
        add_text(s, Inches(0.85), y + Inches(0.38), Inches(11.7),
                 Inches(0.60), body, size=11, color=INK)
        y += Inches(1.05)

    # Money line
    add_rect(s, Inches(0.6), y + Inches(0.05), Inches(12.2), Inches(0.65), NAVY)
    add_text(s, Inches(0.85), y + Inches(0.08), Inches(11.7), Inches(0.60),
             "\"We're not asking for budget to write more skills. "
             "We're asking for the green light to build the system that writes them — "
             "and the domain owners to staff its accuracy gate.\"",
             size=13, bold=True, italic=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)


def build_slide_13_appendix(prs, page, total):
    s = add_blank_slide(prs)
    slide_header(s, "Appendix — sources & artefacts",
                 eyebrow="Backup",
                 page_num=page, total=total)
    sources = [
        ("Article (source of comparison)",
         "How Anthropic enables self-service data analytics with Claude — claude.com/blog, 2026-06-03"),
        ("Internal briefs",
         "proposals/skill-curation-from-nl-and-queries-2026-05-31.md · proposals/skills-mcp-protocol-parity-implementation-2026-06-03.md · dab/monitoring-genie-logs/"),
        ("Skill corpus (knowledge router)",
         "databricks/data-skills/skills/domain-*/ on DataPlatform (CI-deployed)"),
        ("Skill corpus (unbook triple)",
         "/.assistant/skills/data-analysis-playbook · /.assistant/skills/data-analysis-patterns · /.assistant/skills/data-analysis-pattern-library on Databricks Assistant"),
        ("Workspace assistant defaults",
         "/.assistant_workspace_instructions.md at Databricks workspace root"),
        ("Telemetry tables",
         "main.config.monitoring_mcp_logs_mcp_gateway · main.de_output_stg.de_output_monitoring_genie_logs_genie_gateway · main.monitoring.genie_audit_events · main.de_output.de_output_genie_code_skill_feedback"),
        ("Feedback skill (DA-80 PR)",
         "knowledge/skills/feedback-command/SKILL.md"),
    ]
    y = Inches(1.4)
    for h, body in sources:
        add_text(s, Inches(0.6), y, Inches(3.6), Inches(0.4), h, size=11,
                 bold=True, color=NAVY)
        add_text(s, Inches(4.3), y, Inches(8.5), Inches(0.65), body, size=11,
                 color=INK)
        y += Inches(0.62)


# ──────────────────────────── main ────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Order: cover + 12 content + 1 appendix
    builders = [
        build_slide_1_cover,                       # 1 cover
        build_slide_2_why,                         # 2
        build_slide_3_framework,                   # 3
        build_slide_4_scorecard,                   # 4
        build_slide_5_right,                       # 5
        build_slide_6_gaps,                        # 6 — 4 gaps
        build_slide_deliberate_non_adoption,       # 7
        build_slide_open_decision,                 # 8
        build_slide_7_ambivalent_and_designed,     # 9
        build_slide_8_evals,                       # 10
        build_slide_9_provenance,                  # 11
        build_slide_machine,                       # 12 — NEW: vision diagram
        build_slide_10_h2_plan,                    # 13 — now machine-components table
        build_slide_11_asks,                       # 14
        build_slide_12_summary,                    # 15
        build_slide_13_appendix,                   # 16
    ]
    total = len(builders)
    for i, build in enumerate(builders, start=1):
        if i == 1:
            build(prs, total)
        else:
            build(prs, i, total)

    out = Path(__file__).parent / "exec-deck-anthropic-self-service-analytics-2026-06-04.pptx"
    try:
        prs.save(out)
    except PermissionError:
        # File is locked (likely open in PowerPoint) — fall back to suffixed file.
        out = out.with_name(out.stem + "_rev5" + out.suffix)
        prs.save(out)
    print(f"Wrote {out} ({out.stat().st_size:,} bytes, {total} slides)")


if __name__ == "__main__":
    main()
