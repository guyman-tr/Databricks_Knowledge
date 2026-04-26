"""Build subaccount-effort-risk-v2.pptx from the markdown content.

Run: python knowledge/business/_build_subaccount_deck.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).with_name("subaccount-effort-risk-v2.pptx")

DARK_BG    = RGBColor(0x14, 0x14, 0x14)
PANEL_BG   = RGBColor(0x1E, 0x1E, 0x1E)
TEXT       = RGBColor(0xEA, 0xEA, 0xEA)
MUTED      = RGBColor(0xA0, 0xA0, 0xA0)
ACCENT     = RGBColor(0x4F, 0xA8, 0xFF)
GREEN      = RGBColor(0x4A, 0xC4, 0x6B)
YELLOW     = RGBColor(0xE0, 0xB3, 0x4A)
RED        = RGBColor(0xE0, 0x5A, 0x5A)
ORANGE     = RGBColor(0xE8, 0x8A, 0x3C)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bg(slide, color=DARK_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    fill(bg, color)
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""  # clear default
    if isinstance(text, str):
        runs = [(text, {})]
    else:
        runs = text
    first = True
    for run_text, opts in runs:
        if first:
            r = p.add_run()
            first = False
        else:
            r = p.add_run()
        r.text = run_text
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.name = opts.get("font", font)
        r.font.color.rgb = opts.get("color", color)
    return tb


def add_pill(slide, x, y, text, color):
    w = Inches(1.6)
    h = Inches(0.34)
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.5
    fill(s, color)
    tf = s.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.name = "Calibri"
    r.font.color.rgb = RGBColor(0x10, 0x10, 0x10)
    return s


def add_panel(slide, x, y, w, h, color=PANEL_BG):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = 0.04
    fill(s, color)
    return s


def add_footer(slide, n, total):
    add_text(slide, Inches(0.4), Inches(7.1), Inches(8), Inches(0.3),
             "Sub-Account Architecture | Effort & Risk | Apr 23 2026",
             size=10, color=MUTED)
    add_text(slide, Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3),
             f"{n} / {total}",
             size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ---- Slide 1: Title ----
def slide_title(total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    # accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.6), Inches(0.12), Inches(2.0))
    fill(bar, ACCENT)
    add_text(s, Inches(0.9), Inches(2.4), Inches(11), Inches(0.6),
             "Sub-Account Architecture", size=44, bold=True)
    add_text(s, Inches(0.9), Inches(3.1), Inches(11), Inches(0.6),
             "Effort & Risk", size=32, bold=True, color=ACCENT)
    add_text(s, Inches(0.9), Inches(4.0), Inches(11), Inches(0.5),
             "Five options. For each: the specific risk it creates, "
             "and the specific labor that risk forces.",
             size=18, color=MUTED)
    add_text(s, Inches(0.9), Inches(6.4), Inches(11), Inches(0.4),
             "Source: SubAccount alternatives.pdf  |  Apr 23 2026",
             size=12, color=MUTED)


# ---- Slide 2: Two failure modes ----
def slide_failure_modes(n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6),
             "The two failure modes that drive every option",
             size=28, bold=True)

    # Panel 1: Fake-user contamination
    add_panel(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.2))
    add_pill(s, Inches(0.7), Inches(1.7), "OPTIONS 1 + 4", YELLOW)
    add_text(s, Inches(0.7), Inches(2.2), Inches(5.6), Inches(0.5),
             "Fake-user contamination", size=22, bold=True, color=YELLOW)
    add_text(s, Inches(0.7), Inches(2.85), Inches(5.6), Inches(2.2),
             "Sub-account looks like a real customer. Must be excluded "
             "from population/FTD/registration KPIs but counted in "
             "revenue/MIMO/AUM.",
             size=15, color=TEXT)
    add_text(s, Inches(0.7), Inches(4.5), Inches(5.6), Inches(0.4),
             "LABOR SURFACE",
             size=11, bold=True, color=MUTED)
    add_text(s, Inches(0.7), Inches(4.85), Inches(5.6), Inches(1.6),
             [("~140 SPs", {"bold": True, "size": 17}),
              (" filter today on ", {"size": 15}),
              ("IsValidCustomer", {"font": "Consolas", "size": 14}),
              (" or ", {"size": 15}),
              ("IsCreditReportValid", {"font": "Consolas", "size": 14}),
              (". Each is a place that must be re-decided.", {"size": 15})],
             size=15, color=TEXT)

    # Panel 2: Join fanout
    add_panel(s, Inches(6.83), Inches(1.5), Inches(6.0), Inches(5.2))
    add_pill(s, Inches(7.03), Inches(1.7), "OPTIONS 2 + 5", ORANGE)
    add_text(s, Inches(7.03), Inches(2.2), Inches(5.6), Inches(0.5),
             "Join fanout", size=22, bold=True, color=ORANGE)
    add_text(s, Inches(7.03), Inches(2.85), Inches(5.6), Inches(2.4),
             "A previously 1:1 join becomes 1:N.",
             size=15, color=TEXT)
    add_text(s, Inches(7.03), Inches(3.3), Inches(5.6), Inches(1.6),
             [("Opt 2", {"bold": True, "size": 14, "color": RED}),
              (" breaks every ", {"size": 14}),
              ("GCID = GCID", {"font": "Consolas", "size": 13}),
              (" join. Identity-layer change, atomic, ~100+ SPs, hits all 13 Priority 99 finance reports.",
               {"size": 14})],
             size=14, color=TEXT)
    add_text(s, Inches(7.03), Inches(4.85), Inches(5.6), Inches(1.8),
             [("Opt 5", {"bold": True, "size": 14, "color": YELLOW}),
              (" only breaks consumers that join on CID ", {"size": 14}),
              ("without aggregation", {"bold": True, "size": 14}),
              (". Aggregating consumers (SUM + GROUP BY CID) absorb it cleanly. Per-table, gradual, audit-gated.",
               {"size": 14})],
             size=14, color=TEXT)

    add_footer(s, n, total)


# ---- Slide 3: Bottom-line table ----
def slide_bottom_line(n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_text(s, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6),
             "Bottom line", size=28, bold=True)

    rows = [
        ("3", "Mirror No-Copy",                "LOW",      "LOW",      "1 dictionary row + 1 API scope change",      GREEN),
        ("4", "Mirror + Copy (bot CID)",       "MEDIUM",   "MEDIUM",   "Validity-filter audit (~140 SPs)",            YELLOW),
        ("1", "New GCID + CID",                "HIGH",     "MED-HIGH", "Per-KPI bifurcation (~140 SPs)",              YELLOW),
        ("5", "SubPortfolio table (gradual)",  "MED-HIGH", "MEDIUM",   "Per-table consumer audit, 1 table at a time", ORANGE),
        ("2", "New CID -> Same GCID",          "CRITICAL", "CRITICAL", "Schema migration of identity layer",          RED),
    ]
    headers = ["#", "Option", "Effort", "Risk", "What the labor actually is"]
    widths  = [Inches(0.6), Inches(3.4), Inches(1.5), Inches(1.6), Inches(5.2)]
    x0 = Inches(0.5)
    y  = Inches(1.4)
    h  = Inches(0.5)

    # Header row
    cx = x0
    for hdr, w in zip(headers, widths):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y, w, h)
        fill(cell, RGBColor(0x2A, 0x2A, 0x2A))
        tf = cell.text_frame
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.08)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = hdr
        r.font.size = Pt(13); r.font.bold = True
        r.font.name = "Calibri"; r.font.color.rgb = MUTED
        cx += w

    # Body rows
    y += h
    row_h = Inches(0.7)
    for row in rows:
        num, opt, eff, risk, labor, tone = row
        cx = x0
        cells = [num, opt, eff, risk, labor]
        for i, (val, w) in enumerate(zip(cells, widths)):
            cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, y, w, row_h)
            fill(cell, PANEL_BG)
            tf = cell.text_frame
            tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.12)
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.size = Pt(14)
            r.font.name = "Calibri"
            if i == 0:
                r.font.bold = True; r.font.color.rgb = ACCENT
                p.alignment = PP_ALIGN.CENTER
            elif i == 1:
                r.font.bold = True; r.font.color.rgb = TEXT
            elif i in (2, 3):
                r.font.bold = True; r.font.color.rgb = tone
            else:
                r.font.color.rgb = TEXT
            cx += w
        y += row_h

    add_footer(s, n, total)


# ---- Per-option slides ----
OPTIONS = [
    {
        "id": 3,
        "name": "Mirror No-Copy",
        "effort": ("LOW", GREEN),
        "risk":   ("LOW", GREEN),
        "headline": "Add a MirrorTypeID. No new identity, no new joins, no fake users.",
        "the_risk": "None of consequence. Same row stays in every consumer.",
        "the_labor": "Add 1 row to Dictionary.MirrorType. Plumb mirrorID into the trading API JWT scope. Decide private-instrument visibility for the sub-account.",
        "extra_label": "Hard ceiling",
        "extra": "No copy support inside the sub-account.",
    },
    {
        "id": 4,
        "name": "Mirror + Copy (bot CID)",
        "effort": ("MEDIUM", YELLOW),
        "risk":   ("MEDIUM", YELLOW),
        "headline": "Real bot accounts created underneath the user; copy engine handles trading.",
        "the_risk": "Bot CID looks like a real customer everywhere. Pollutes user counts, FTD, registrations, AML alerts, KYC pipeline. Compliance pattern flag: a customer with zero organic activity that only ever copies one CID is exactly what AML alerts trigger on.",
        "the_labor": "Add IsBotAccount (or extend IsInternal) at the OLTP customer record. Then add it to the WHERE clause of every analytics SP currently filtering by IsValidCustomer and IsCreditReportValid - ~140 SPs. AML/KYC pipelines need explicit review (they may want to keep watching bot CIDs even when other reports exclude them).",
    },
    {
        "id": 1,
        "name": "New GCID + CID",
        "effort": ("HIGH", YELLOW),
        "risk":   ("MED-HIGH", YELLOW),
        "headline": "Sub-account is a fully separate customer. User logs in twice.",
        "the_risk": "Same fake-user contamination as Option 4, plus acquisition KPIs (FTD, Registrations, valid-customer counts) inflate per sub-account. Tax filings (1099, W8Ben) emit per GCID, so a sub-GCID = duplicate IRS filing risk if not excluded.",
        "the_labor": "Same IsBotAccount-style plumbing as Option 4, but you also decide PER-METRIC whether the sub-GCID is counted (revenue/MIMO/AUM = yes; FTD/Registration/valid-customer = no). Two deploys, not one. Same ~140 SP review surface, with bigger semantic decisions per SP.",
    },
    {
        "id": 5,
        "name": "SubPortfolio table (gradual)",
        "effort": ("MED-HIGH", ORANGE),
        "risk":   ("MEDIUM", YELLOW),
        "headline": "New (CID, SubPortfolioID) grain on selected tables. No fake users, no PK changes.",
        "the_risk": "Per broken table, any consumer that joins on CID without subsequent aggregation gets row fanout. Aggregating consumers (SUM + GROUP BY CID downstream) absorb it cleanly. Probably not huge but not currently quantified.",
        "the_labor": "Per table you choose to break: classify every direct consumer as (a) aggregator -> safe, (b) pass-through-by-CID -> must be updated. One table at a time, audit-gated. Save Priority 99 finance SPs for last.",
        "extra_label": "Gate before committing",
        "extra": "Half-day static SQL scan: per producer table, list direct consumers and classify each as safe (has SUM + GROUP BY CID downstream) or unsafe (selects CID rows without aggregation). Turns Option 5 risk into a sized backlog.",
    },
    {
        "id": 2,
        "name": "New CID -> Same GCID",
        "effort": ("CRITICAL", RED),
        "risk":   ("CRITICAL", RED),
        "headline": "Multiple CIDs share one GCID. Breaks the GCID PK on Customer.CustomerIdentification.",
        "the_risk": "Every JOIN ON x.GCID = y.GCID returns N rows where 1 was expected. Hundreds of SPs do this. Identity-layer change, atomic, unrecoverable. Direct hit to all 13 Priority 99 finance reports.",
        "the_labor": "Migrate Customer.CustomerIdentification PK. Update every GCID-keyed JOIN in BI_DB / DWH / Dealing / eMoney / EXW + 5,000+ ComplianceDB objects + every OLTP source DB in lockstep. No partial deploy.",
    },
]


def slide_option(opt, n, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)

    # Header: "Option N - Name"
    add_text(s, Inches(0.5), Inches(0.4), Inches(8.5), Inches(0.7),
             [(f"Option {opt['id']}  ", {"size": 18, "color": MUTED, "bold": True}),
              (opt["name"], {"size": 28, "bold": True, "color": TEXT})])

    # Effort/Risk pills (top-right)
    add_pill(s, Inches(10.1), Inches(0.5), f"Effort {opt['effort'][0]}", opt["effort"][1])
    add_pill(s, Inches(11.75), Inches(0.5), f"Risk {opt['risk'][0]}",   opt["risk"][1])

    # Headline panel
    add_panel(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.85),
              color=RGBColor(0x22, 0x2A, 0x36))
    add_text(s, Inches(0.75), Inches(1.5), Inches(11.8), Inches(0.7),
             opt["headline"], size=17, bold=True, color=ACCENT)

    # The risk
    y = Inches(2.55)
    add_text(s, Inches(0.5), y, Inches(12.3), Inches(0.4),
             "THE RISK", size=12, bold=True, color=RED)
    add_text(s, Inches(0.5), Inches(2.95), Inches(12.3), Inches(1.4),
             opt["the_risk"], size=15, color=TEXT)

    # The labor
    y2 = Inches(4.4)
    add_text(s, Inches(0.5), y2, Inches(12.3), Inches(0.4),
             "THE LABOR IT CREATES", size=12, bold=True, color=YELLOW)
    add_text(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.6),
             opt["the_labor"], size=15, color=TEXT)

    # Optional extra
    if "extra" in opt:
        add_panel(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7),
                  color=RGBColor(0x2A, 0x24, 0x18))
        add_text(s, Inches(0.7), Inches(6.35), Inches(2.5), Inches(0.3),
                 opt["extra_label"].upper(), size=11, bold=True, color=ORANGE)
        add_text(s, Inches(0.7), Inches(6.62), Inches(11.9), Inches(0.4),
                 opt["extra"], size=12, color=TEXT)

    add_footer(s, n, total)


total = 2 + 1 + len(OPTIONS)  # title + failure-modes + bottom-line + 5 options
slide_title(total)
slide_failure_modes(2, total)
slide_bottom_line(3, total)
for i, opt in enumerate(OPTIONS):
    slide_option(opt, 4 + i, total)

prs.save(OUT)
print(f"Wrote {OUT}")
